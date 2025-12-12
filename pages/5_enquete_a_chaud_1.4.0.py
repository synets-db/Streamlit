import streamlit as st
import pandas as pd
import altair as alt
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt

# ==============================================================
# Helpers : Altair -> PNG (bytes) via vl-convert
# ==============================================================
def chart_to_png_bytes(chart: alt.Chart, scale: float = 2.0) -> bytes:
    """
    Convertit un chart Altair (Vega-Lite) en PNG bytes via vl-convert-python.
    """
    import vl_convert as vlc  # pip install vl-convert-python
    spec = chart.to_dict()
    return vlc.vegalite_to_png(spec, scale=scale)


# ==============================================================
# PPTX : génération avec camemberts
# ==============================================================
def generate_pptx_with_charts(items_data, open_questions, title_main="Enquête à chaud"):
    prs = Presentation()

    # --- Slide titre ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_main
    slide.placeholders[1].text = "Rapport d'analyse"

    blank = prs.slide_layouts[6]  # blank slide

    for it in items_data:
        slide = prs.slides.add_slide(blank)

        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
        tf = title_box.text_frame
        tf.text = it["item"]
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True

        # Sous-titre total
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(6.5), Inches(0.35))
        sub_tf = sub_box.text_frame
        sub_tf.text = f"Total réponses : {it['total']}" if it["total"] else "Total réponses : N/A"
        sub_tf.paragraphs[0].font.size = Pt(12)

        # Camembert (gauche)
        if it.get("pie") is not None and it.get("total", 0) and it["total"] > 0:
            try:
                png_bytes = chart_to_png_bytes(it["pie"], scale=2.0)
                img_stream = BytesIO(png_bytes)
                # Ajuste la taille/position selon ton goût
                slide.shapes.add_picture(img_stream, Inches(0.6), Inches(1.4), width=Inches(5.2))
            except Exception:
                # si conversion KO, on continue sans image
                pass

        # Détails (droite)
        details_box = slide.shapes.add_textbox(Inches(6.0), Inches(1.4), Inches(6.8), Inches(3.2))
        dtf = details_box.text_frame
        dtf.word_wrap = True
        dtf.clear()

        if it.get("modalites"):
            p0 = dtf.paragraphs[0]
            p0.text = "Résultats :"
            p0.font.bold = True
            p0.font.size = Pt(14)

            for row in it["modalites"]:
                mod = row.get("Modalité", "")
                nb = row.get("Nombre", 0)
                pct = row.get("Pourcentage", 0)
                p = dtf.add_paragraph()
                p.text = f"{mod} : {nb} ({pct}%)"
                p.level = 1
                p.font.size = Pt(12)
        else:
            p0 = dtf.paragraphs[0]
            p0.text = "Aucun résultat."
            p0.font.size = Pt(12)

        # Commentaires (bas)
        coms = it.get("commentaires", [])
        if coms:
            com_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.85), Inches(12.2), Inches(2.0))
            ctf = com_box.text_frame
            ctf.word_wrap = True
            ctf.text = "Commentaires :"
            ctf.paragraphs[0].font.bold = True
            ctf.paragraphs[0].font.size = Pt(14)

            for txt in coms:
                p = ctf.add_paragraph()
                p.text = f"- {txt}"
                p.level = 1
                p.font.size = Pt(11)

    # Slides questions ouvertes finales
    for q in open_questions:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
        slide.shapes.title.text = q["question"]
        tf = slide.placeholders[1].text_frame
        tf.clear()

        if not q["reponses"]:
            tf.text = "Aucune réponse."
        else:
            tf.text = "Réponses :"
            for rep in q["reponses"]:
                p = tf.add_paragraph()
                p.text = f"- {rep}"
                p.level = 1

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# ==============================================================
# Streamlit page
# ==============================================================
st.set_page_config(page_title="Enquête à chaud — Items", layout="wide")
st.title("📝 Enquête à chaud — Analyse de tous les items")

st.markdown("""
- Analyse des **items C → AZ** : camembert + détail + commentaires (en **2 colonnes**).
- Puis affichage des **questions ouvertes finales** (BO / BQ / BS).
- Export **PowerPoint** avec **camemberts intégrés**.
""")

uploaded_file = st.file_uploader(
    "Déposez le fichier CSV (export Evento, séparateur `;`)",
    type=["csv"]
)

if not uploaded_file:
    st.info("En attente d'un fichier CSV…")
    st.stop()

# Lecture CSV
try:
    df_raw = pd.read_csv(uploaded_file, sep=";", encoding="utf-8")
except UnicodeDecodeError:
    df_raw = pd.read_csv(uploaded_file, sep=";", encoding="latin-1")

if df_raw.shape[1] <= 2:
    st.error("Le fichier ne contient pas suffisamment de colonnes (A et B + réponses).")
    st.stop()

# Colonnes C → AZ = indices 2..51 (slice 2:52)
max_col = min(df_raw.shape[1], 52)
df = df_raw.iloc[:, 2:max_col].copy()

st.write(f"Nombre de lignes : {len(df)}")
st.write(f"Colonnes analysées (C → AZ) : {df.shape[1]}")

last_row = df.iloc[-1]
df_without_total = df.iloc[:-1]

# items = { item_label: {"mod_cols": [(modalité, col)], "comment_cols": [col]} }
items = {}
for col in df.columns:
    full = str(col)
    if "." in full:
        base, modality = full.rsplit(".", 1)
    else:
        base, modality = full, ""
    base = base.strip()
    modality = modality.strip()

    items.setdefault(base, {"mod_cols": [], "comment_cols": []})

    if "commentaire" in modality.lower():
        items[base]["comment_cols"].append(col)
    else:
        items[base]["mod_cols"].append((modality, col))

# Palette satisfaction
ordre_satisfaction = ["Très satisfait", "Satisfait", "Peu satisfait", "Pas du tout satisfait"]
couleurs_satisfaction = {
    "Très satisfait": "#1b7837",
    "Satisfait": "#5aae61",
    "Peu satisfait": "#80cdc1",
    "Pas du tout satisfait": "#f46d43",
}

# Structures export PPTX
export_items = []
export_open_questions = []

# ==============================================================
# Affichage items en 2 colonnes
# ==============================================================
for item_label, info in items.items():
    mod_cols = info["mod_cols"]
    comment_cols = info["comment_cols"]

    if not mod_cols and not comment_cols:
        continue

    chart_df = None
    pie = None
    total = 0

    # --- modalités (totaux sur dernière ligne) ---
    if mod_cols:
        labels, counts = [], []
        for modality, col in mod_cols:
            labels.append(modality)
            try:
                val = int(last_row[col])
            except (ValueError, TypeError):
                val = 0
            counts.append(val)

        total = sum(counts)
        if total > 0:
            percentages = [round(c / total * 100, 1) for c in counts]
            chart_df = pd.DataFrame({"Modalité": labels, "Nombre": counts, "Pourcentage": percentages})

            # satisfaction 4 niveaux
            if set(labels).issubset(set(ordre_satisfaction)):
                chart_df["Modalité"] = pd.Categorical(
                    chart_df["Modalité"],
                    categories=ordre_satisfaction,
                    ordered=True
                )
                chart_df = chart_df.sort_values("Modalité")

                pie = (
                    alt.Chart(chart_df)
                    .mark_arc(innerRadius=40)
                    .encode(
                        theta="Nombre:Q",
                        color=alt.Color(
                            "Modalité:N",
                            scale=alt.Scale(
                                domain=ordre_satisfaction,
                                range=[couleurs_satisfaction[m] for m in ordre_satisfaction]
                            ),
                            legend=alt.Legend(title="Modalité")
                        ),
                        tooltip=["Modalité", "Nombre", "Pourcentage"]
                    )
                    .properties(width=400, height=200)  # <- ton réglage
                )
            else:
                pie = (
                    alt.Chart(chart_df)
                    .mark_arc(innerRadius=40)
                    .encode(
                        theta="Nombre:Q",
                        color=alt.Color("Modalité:N", legend=alt.Legend(title="Modalité")),
                        tooltip=["Modalité", "Nombre", "Pourcentage"]
                    )
                    .properties(width=400, height=200)  # <- ton réglage
                )

    # --- commentaires (toutes lignes sauf dernière) ---
    commentaires = []
    if comment_cols:
        for col in comment_cols:
            s = df_without_total[col].dropna().astype(str).str.strip()
            s = s[s != ""]
            commentaires.extend(list(s))

    # --- UI ---
    st.markdown(f"## ❓ {item_label}")
    col1, col2 = st.columns([1, 1])

    with col1:
        if mod_cols and total > 0 and pie is not None:
            st.markdown(f"**{total} réponses**")
            st.altair_chart(pie, use_container_width=False)
        elif mod_cols:
            st.info("Aucun total disponible pour calculer les pourcentages.")
        else:
            st.info("Aucune modalité fermée (uniquement commentaire).")

    with col2:
        if mod_cols and chart_df is not None and total > 0:
            st.markdown("### 📊 Détail des réponses")
            for lab, c, p in zip(chart_df["Modalité"], chart_df["Nombre"], chart_df["Pourcentage"]):
                st.markdown(f"- **{lab}** : {c} réponses ({p}%)")

        if comment_cols:
            st.markdown("### 💬 Commentaires")
            if not commentaires:
                st.info("Aucun commentaire.")
            else:
                for i, txt in enumerate(commentaires, start=1):
                    st.markdown(f"- **Commentaire {i}** : {txt}")

    st.markdown("---")

    # --- export data ---
    export_items.append({
        "item": item_label,
        "total": total,
        "modalites": chart_df.to_dict("records") if chart_df is not None else [],
        "commentaires": commentaires,
        "pie": pie,  # ✅ important pour PPTX avec camemberts
    })

# ==============================================================
# Questions ouvertes finales : BO, BQ, BS
# ==============================================================
st.header("📝 Questions ouvertes finales")

open_indices = [66, 68, 70]  # BO, BQ, BS (0-based)
ncols = df_raw.shape[1]

for idx in open_indices:
    if idx >= ncols:
        continue

    col_name = str(df_raw.columns[idx])
    if ".Commentaire" in col_name:
        question_label = col_name.replace(".Commentaire", "").strip()
    else:
        question_label = col_name.split(".")[0].strip()

    st.subheader(f"❓ {question_label}")

    serie = df_raw.iloc[:-1, idx].dropna().astype(str).str.strip()
    serie = serie[serie != ""]
    reponses = list(serie)

    if not reponses:
        st.info("Aucune réponse.")
    else:
        for i, txt in enumerate(reponses, start=1):
            st.markdown(f"- **Réponse {i}** : {txt}")

    export_open_questions.append({"question": question_label, "reponses": reponses})

# ==============================================================
# Export PPTX (avec camemberts)
# ==============================================================
st.divider()
st.subheader("📊 Export du rapport")

pptx_buffer = generate_pptx_with_charts(
    export_items,
    export_open_questions,
    title_main="Enquête à chaud"
)

st.download_button(
    label="📊 Télécharger le rapport PowerPoint (.pptx)",
    data=pptx_buffer,
    file_name="enquete_a_chaud_rapport.pptx",
    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
)
